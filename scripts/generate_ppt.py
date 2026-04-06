#!/usr/bin/env python3
"""Generate PowerPoint presentations from a JSON slide definition.

Usage:
    python generate_ppt.py --slides slides.json --output deck.pptx [--title "Title"] [--widescreen]

Slide JSON format: array of slide objects, each with a "layout" key and layout-specific fields.
See ~/.claude/commands/generate-ppt.md for the full schema.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --- Modern professional theme ---
COLOR_TITLE = RGBColor(0x1A, 0x1A, 0x2E)       # Deep navy for titles
COLOR_BODY = RGBColor(0x2D, 0x2D, 0x2D)         # Charcoal for body text
COLOR_ACCENT = RGBColor(0x00, 0x6D, 0x77)       # Teal accent (modern, professional)
COLOR_ACCENT_LIGHT = RGBColor(0xE0, 0xF2, 0xF1) # Light teal background
COLOR_LIGHT = RGBColor(0x5A, 0x5A, 0x6E)        # Muted gray-blue for subtitles
COLOR_TABLE_HEADER_BG = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TABLE_ALT_BG = RGBColor(0xF5, 0xF5, 0xF7)
COLOR_DIVIDER = RGBColor(0xE0, 0xE0, 0xE0)

# --- Font defaults (modern sans-serif) ---
FONT_TITLE = "Segoe UI"
FONT_BODY = "Segoe UI"
FONT_SIZE_TITLE = Pt(36)
FONT_SIZE_SUBTITLE = Pt(18)
FONT_SIZE_SECTION = Pt(34)
FONT_SIZE_SLIDE_TITLE = Pt(26)
FONT_SIZE_BODY = Pt(17)
FONT_SIZE_BODY_SUB = Pt(15)
FONT_SIZE_TABLE = Pt(13)
FONT_SIZE_TABLE_HEADER = Pt(14)
FONT_SIZE_NOTES = Pt(12)


def create_presentation(widescreen: bool = True) -> Presentation:
    prs = Presentation()
    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs


def _style_run(run, font_name: str, font_size, color: RGBColor, bold: bool = False):
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold


def _add_textbox(slide, left, top, width, height, text: str, font_name: str = FONT_BODY,
                 font_size=FONT_SIZE_BODY, color: RGBColor = COLOR_BODY,
                 bold: bool = False, alignment=PP_ALIGN.LEFT, word_wrap: bool = True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _style_run(run, font_name, font_size, color, bold)
    return txBox


def _add_bullets(text_frame, bullets: list[str], font_size=FONT_SIZE_BODY,
                 sub_font_size=FONT_SIZE_BODY_SUB, color: RGBColor = COLOR_BODY):
    """Add bullet points to a text frame. Supports nested bullets via leading '- '."""
    text_frame.clear()
    for i, bullet in enumerate(bullets):
        level = 0
        text = bullet
        # Support simple nesting: items starting with "  " or "\t" are sub-bullets
        while text.startswith("  ") or text.startswith("\t"):
            level += 1
            text = text.lstrip("\t")
            if text.startswith("  "):
                text = text[2:]
        # Strip leading "- " or "* " markers
        for marker in ("- ", "* ", "-> "):
            if text.startswith(marker):
                text = text[len(marker):]
                break

        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()

        p.level = level
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = text
        size = sub_font_size if level > 0 else font_size
        _style_run(run, FONT_BODY, size, color)


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    sw = prs.slide_width
    sh = prs.slide_height

    # Title - centered in upper portion
    _add_textbox(slide, Inches(1), Inches(2), sw - Inches(2), Inches(1.5),
                 title, FONT_TITLE, FONT_SIZE_TITLE, COLOR_TITLE, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    if subtitle:
        _add_textbox(slide, Inches(1.5), Inches(3.8), sw - Inches(3), Inches(1),
                     subtitle, FONT_BODY, FONT_SIZE_SUBTITLE, COLOR_LIGHT,
                     alignment=PP_ALIGN.CENTER)

    # Accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.5),
        sw - Inches(9), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    return slide


def add_section_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sw = prs.slide_width
    sh = prs.slide_height

    # Section title - centered
    _add_textbox(slide, Inches(1), Inches(2.5), sw - Inches(2), Inches(2),
                 title, FONT_TITLE, FONT_SIZE_SECTION, COLOR_ACCENT, bold=True,
                 alignment=PP_ALIGN.CENTER)

    # Accent line

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.7),
        sw - Inches(10), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    return slide


def add_content_slide(prs: Presentation, title: str, bullets: list[str],
                      notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sw = prs.slide_width

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.4), sw - Inches(1.6), Inches(1),
                 title, FONT_TITLE, FONT_SIZE_SLIDE_TITLE, COLOR_TITLE, bold=True)

    # Accent underline

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35),
        Inches(2), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Bullets
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(1.7), sw - Inches(2), Inches(5.2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    _add_bullets(tf, bullets)

    # Speaker notes
    if notes:
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_tf.text = notes

    return slide


def add_two_column_slide(prs: Presentation, title: str,
                         left_title: str, left_bullets: list[str],
                         right_title: str, right_bullets: list[str],
                         notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sw = prs.slide_width
    col_width = (sw - Inches(2.2)) / 2

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.4), sw - Inches(1.6), Inches(1),
                 title, FONT_TITLE, FONT_SIZE_SLIDE_TITLE, COLOR_TITLE, bold=True)

    # Accent underline

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.35),
        Inches(2), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Left column title
    _add_textbox(slide, Inches(0.8), Inches(1.6), col_width, Inches(0.6),
                 left_title, FONT_TITLE, FONT_SIZE_SUBTITLE, COLOR_ACCENT, bold=True)

    # Left column bullets
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.3), col_width - Inches(0.4), Inches(4.5))
    txBox.text_frame.word_wrap = True
    _add_bullets(txBox.text_frame, left_bullets)

    # Right column title
    right_left = Inches(0.8) + col_width + Inches(0.6)
    _add_textbox(slide, right_left, Inches(1.6), col_width, Inches(0.6),
                 right_title, FONT_TITLE, FONT_SIZE_SUBTITLE, COLOR_ACCENT, bold=True)

    # Right column bullets
    txBox = slide.shapes.add_textbox(right_left + Inches(0.2), Inches(2.3),
                                     col_width - Inches(0.4), Inches(4.5))
    txBox.text_frame.word_wrap = True
    _add_bullets(txBox.text_frame, right_bullets)

    # Vertical divider
    divider_left = Inches(0.8) + col_width + Inches(0.15)
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, divider_left, Inches(1.7),
        Inches(0.02), Inches(5)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = COLOR_DIVIDER
    div.line.fill.background()

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_table_slide(prs: Presentation, title: str, headers: list[str],
                    rows: list[list[str]], notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sw = prs.slide_width

    # Title
    _add_textbox(slide, Inches(0.8), Inches(0.4), sw - Inches(1.6), Inches(1),
                 title, FONT_TITLE, FONT_SIZE_SLIDE_TITLE, COLOR_TITLE, bold=True)

    # Table
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    table_width = sw - Inches(2)
    table_height = Inches(min(num_rows * 0.5, 5.5))

    graphic_frame = slide.shapes.add_table(
        num_rows, num_cols, Inches(1), Inches(1.6), table_width, table_height
    )
    table = graphic_frame.table

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = header
        _style_run(run, FONT_BODY, FONT_SIZE_TABLE_HEADER, COLOR_TABLE_HEADER_FG, bold=True)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER_BG
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            _style_run(run, FONT_BODY, FONT_SIZE_TABLE, COLOR_BODY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternating row colors
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_TABLE_ALT_BG

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


def add_blank_slide(prs: Presentation, title: str = "", text: str = "",
                    notes: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    sw = prs.slide_width
    sh = prs.slide_height

    if title:
        _add_textbox(slide, Inches(0.8), Inches(0.4), sw - Inches(1.6), Inches(1),
                     title, FONT_TITLE, FONT_SIZE_SLIDE_TITLE, COLOR_TITLE, bold=True)

    if text:
        _add_textbox(slide, Inches(2), Inches(2.5), sw - Inches(4), Inches(3),
                     text, FONT_BODY, FONT_SIZE_SUBTITLE, COLOR_BODY,
                     alignment=PP_ALIGN.CENTER)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes

    return slide


# --- Dispatch ---

BUILDERS = {
    "title": lambda prs, s: add_title_slide(prs, s["title"], s.get("subtitle", "")),
    "section": lambda prs, s: add_section_slide(prs, s["title"]),
    "content": lambda prs, s: add_content_slide(
        prs, s["title"], s.get("bullets", []), s.get("notes", "")),
    "two_column": lambda prs, s: add_two_column_slide(
        prs, s["title"],
        s.get("left_title", ""), s.get("left_bullets", []),
        s.get("right_title", ""), s.get("right_bullets", []),
        s.get("notes", "")),
    "table": lambda prs, s: add_table_slide(
        prs, s["title"], s.get("headers", []), s.get("rows", []), s.get("notes", "")),
    "blank": lambda prs, s: add_blank_slide(
        prs, s.get("title", ""), s.get("text", ""), s.get("notes", "")),
}


def build_presentation(slides_data: list[dict], widescreen: bool = True) -> Presentation:
    prs = create_presentation(widescreen)

    for i, slide_def in enumerate(slides_data):
        layout = slide_def.get("layout", "content")
        builder = BUILDERS.get(layout)
        if not builder:
            print(f"Warning: Unknown layout '{layout}' on slide {i+1}, using 'content'",
                  file=sys.stderr)
            builder = BUILDERS["content"]
        builder(prs, slide_def)

    return prs


def main():
    parser = argparse.ArgumentParser(description="Generate a PowerPoint from JSON slide definitions")
    parser.add_argument("--slides", required=True, help="Path to JSON file with slide definitions")
    parser.add_argument("--output", required=True, help="Output .pptx file path")
    parser.add_argument("--widescreen", action="store_true", default=True,
                        help="Use 16:9 widescreen (default)")
    parser.add_argument("--standard", action="store_true",
                        help="Use 4:3 standard aspect ratio")
    args = parser.parse_args()

    slides_path = Path(args.slides)
    if not slides_path.exists():
        print(f"Error: Slides file not found: {slides_path}", file=sys.stderr)
        sys.exit(1)

    with open(slides_path, encoding="utf-8") as f:
        slides_data = json.load(f)

    widescreen = not args.standard
    prs = build_presentation(slides_data, widescreen=widescreen)

    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    print(f"Generated {len(slides_data)} slides -> {output_path}")


if __name__ == "__main__":
    main()
